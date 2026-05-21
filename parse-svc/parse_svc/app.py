"""Document parsing service — PDF, DOCX, PPTX, XLSX, and more to markdown.

Routes files to the right parser based on extension, with OCR fallback
for scanned PDFs and table extraction for tabular data.
"""

import io
import logging
import os
import re
import tempfile
from pathlib import Path
from typing import Any

from fastapi import FastAPI, HTTPException, UploadFile
from pydantic import BaseModel

logger = logging.getLogger(__name__)

app = FastAPI(title="GroktoCrawl Parse Service", version="0.1.0")

MAX_SIZE_MB = int(os.getenv("PARSE_MAX_SIZE_MB", "50"))
MAX_SIZE_BYTES = MAX_SIZE_MB * 1024 * 1024


class ParseResponse(BaseModel):
    success: bool
    data: dict[str, Any] | None = None
    error: str | None = None


# ---- Format detection ----

def _ext(filename: str) -> str:
    return Path(filename).suffix.lower().lstrip(".")


SUPPORTED_FORMATS = {
    "pdf": "PDF document",
    "docx": "Word document",
    "pptx": "PowerPoint presentation",
    "xlsx": "Excel workbook",
    "csv": "CSV data",
    "md": "Markdown",
    "txt": "Plain text",
    "json": "JSON data",
    "yaml": "YAML data",
    "yml": "YAML data",
    "xml": "XML data",
    "html": "HTML document",
    "htm": "HTML document",
}


# ---- Individual parsers ----


def _parse_pdf(content: bytes, filename: str) -> dict:
    """Parse PDF — try text extraction first, fall back to OCR."""
    markdown_parts = []
    metadata: dict[str, Any] = {"format": "pdf", "filename": filename}
    text_extracted = False

    # Tier 1: pypdf text extraction
    try:
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(content))
        metadata["pages"] = len(reader.pages)
        pages_text = []
        for page in reader.pages:
            text = page.extract_text() or ""
            pages_text.append(text)
        full_text = "\n\n".join(pages_text).strip()
        if len(full_text) > 50:
            markdown_parts.append(full_text)
            text_extracted = True
            metadata["extraction"] = "pypdf"
    except Exception as e:
        logger.warning("pypdf failed: %s", e)

    # Tier 2: OCR for scanned PDFs (if pypdf got little or nothing)
    if not text_extracted or len("".join(markdown_parts)) < 100:
        try:
            import pytesseract
            from pdf2image import convert_from_bytes

            images = convert_from_bytes(content, dpi=300, first_page=1, last_page=min(10, metadata.get("pages", 10)))
            ocr_parts = []
            for i, img in enumerate(images):
                text = pytesseract.image_to_string(img, lang="eng")
                ocr_parts.append(f"--- Page {i+1} ---\n\n{text}")
            if ocr_parts:
                ocr_text = "\n\n".join(ocr_parts).strip()
                if len(ocr_text) > 50:
                    markdown_parts = [ocr_text]
                    metadata["extraction"] = "ocr"
                    text_extracted = True
        except ImportError:
            logger.warning("pytesseract/pdf2image not available, skipping OCR")
        except Exception as e:
            logger.warning("OCR failed: %s", e)

    # Tier 3: Table extraction for PDFs with tables
    try:
        import camelot

        tables = camelot.read_pdf(io.BytesIO(content), pages="all", flavor="lattice")
        if len(tables) > 0:
            from tabulate import tabulate

            table_md_parts = []
            for i, table in enumerate(tables):
                md = tabulate(table.df, headers="keys", tablefmt="github")
                table_md_parts.append(md)
            if table_md_parts:
                markdown_parts.append("\n\n### Extracted Tables\n\n" + "\n\n".join(table_md_parts))
                metadata["tables_found"] = len(tables)
    except ImportError:
        logger.debug("camelot not available, skipping table extraction")
    except Exception as e:
        logger.debug("Table extraction failed: %s", e)

    return {"markdown": "\n\n".join(markdown_parts).strip(), "metadata": metadata}


def _parse_docx(content: bytes, filename: str) -> dict:
    """Parse DOCX to markdown."""
    try:
        from docx import Document

        doc = Document(io.BytesIO(content))
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        text = "\n\n".join(paragraphs)

        # Extract tables
        table_parts = []
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(cells)
            if rows:
                from tabulate import tabulate

                table_parts.append(tabulate(rows, headers="firstrow", tablefmt="github"))
        if table_parts:
            text += "\n\n### Tables\n\n" + "\n\n".join(table_parts)

        metadata = {
            "format": "docx",
            "filename": filename,
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables),
        }
        return {"markdown": text.strip(), "metadata": metadata}
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Failed to parse DOCX: {e}")


def _parse_pptx(content: bytes, filename: str) -> dict:
    """Parse PPTX to markdown — slides, notes, and speaker notes."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"## Slide {i}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_parts.append(para.text)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        from tabulate import tabulate
                        slide_parts.append(tabulate(rows, headers="firstrow", tablefmt="github"))
            slides_text.append("\n\n".join(slide_parts))

        metadata = {
            "format": "pptx",
            "filename": filename,
            "slides": len(prs.slides),
        }
        return {"markdown": "\n\n---\n\n".join(slides_text).strip(), "metadata": metadata}
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Failed to parse PPTX: {e}")


def _parse_xlsx(content: bytes, filename: str) -> dict:
    """Parse XLSX to markdown tables."""
    try:
        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        sheets_md = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append([str(c) if c is not None else "" for c in row])
            if rows:
                from tabulate import tabulate
                sheets_md.append(f"### Sheet: {sheet_name}\n\n{tabulate(rows, headers='firstrow', tablefmt='github')}")

        wb.close()
        metadata = {
            "format": "xlsx",
            "filename": filename,
            "sheets": len(wb.sheetnames),
        }
        return {"markdown": "\n\n".join(sheets_md).strip(), "metadata": metadata}
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Failed to parse XLSX: {e}")


def _parse_text(content: bytes, filename: str) -> dict:
    """Plain text / code / markdown files."""
    text = content.decode("utf-8", errors="replace")
    metadata = {"format": _ext(filename), "filename": filename, "chars": len(text)}
    return {"markdown": text.strip(), "metadata": metadata}


def _parse_csv(content: bytes, filename: str) -> dict:
    """Parse CSV to markdown table."""
    text = content.decode("utf-8", errors="replace")
    import csv
    import io as io_module

    reader = csv.reader(io_module.StringIO(text))
    rows = list(reader)
    if rows:
        from tabulate import tabulate
        md = tabulate(rows, headers="firstrow", tablefmt="github")
    else:
        md = text
    metadata = {"format": "csv", "filename": filename, "rows": len(rows)}
    return {"markdown": md, "metadata": metadata}


def _parse_via_markitdown(content: bytes, filename: str) -> dict:
    """Use Microsoft markitdown as a catch-all fallback."""
    try:
        from markitdown import MarkItDown

        md_converter = MarkItDown()
        with tempfile.NamedTemporaryFile(suffix=f".{_ext(filename)}", delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        try:
            result = md_converter.convert(tmp_path)
            return {"markdown": result.text_content.strip(), "metadata": {"format": _ext(filename), "filename": filename, "extraction": "markitdown"}}
        finally:
            os.unlink(tmp_path)
    except Exception as e:
        logger.warning("markitdown failed for %s: %s", filename, e)
        return {"markdown": content.decode("utf-8", errors="replace")[:50000], "metadata": {"format": _ext(filename), "filename": filename, "extraction": "raw"}}


# ---- Router ----

PARSERS = {
    "pdf": _parse_pdf,
    "docx": _parse_docx,
    "pptx": _parse_pptx,
    "xlsx": _parse_xlsx,
    "csv": _parse_csv,
    "md": _parse_text,
    "txt": _parse_text,
    "json": _parse_text,
    "yaml": _parse_text,
    "yml": _parse_text,
    "xml": _parse_text,
    "html": _parse_text,
    "htm": _parse_text,
}


@app.get("/health")
async def health():
    return {"status": "ok"}


@app.post("/parse", response_model=ParseResponse)
async def parse_file(file: UploadFile):
    """Upload a file and get its content as markdown."""
    if not file.filename:
        raise HTTPException(status_code=400, detail="No filename provided")

    ext = _ext(file.filename)
    if not ext:
        raise HTTPException(status_code=400, detail=f"Could not determine file extension: {file.filename}")

    if ext not in SUPPORTED_FORMATS:
        raise HTTPException(status_code=400, detail=f"Unsupported format: .{ext}. Supported: {', '.join(sorted(SUPPORTED_FORMATS))}")

    content = await file.read()
    if len(content) > MAX_SIZE_BYTES:
        raise HTTPException(status_code=413, detail=f"File too large. Max {MAX_SIZE_MB}MB.")

    logger.info("Parsing %s (%s, %d bytes)", file.filename, ext, len(content))

    parser = PARSERS.get(ext, _parse_via_markitdown)
    try:
        result = parser(content, file.filename)
        md = result.get("markdown", "")
        meta = result.get("metadata", {})
        meta["size_bytes"] = len(content)
        return ParseResponse(success=True, data={"markdown": md, "metadata": meta})
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("Parse failed for %s", file.filename)
        return ParseResponse(success=False, error=str(e))
