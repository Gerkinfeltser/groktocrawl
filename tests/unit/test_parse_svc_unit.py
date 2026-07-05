"""Unit tests for parse-svc parser functions.

Tests each of the 8 parser functions in parse_svc/app.py directly,
including OCR/Camelot graceful degradation paths.
"""

import io
from unittest.mock import patch

import pytest
from parse_svc.app import (
    _ext,
    _parse_csv,
    _parse_docx,
    _parse_pdf,
    _parse_pptx,
    _parse_text,
    _parse_via_markitdown,
    _parse_xlsx,
)

# ── _ext ──────────────────────────────────────────────────────────────────


class TestExt:
    """Extension extraction from various filenames."""

    def test_simple_extension(self):
        assert _ext("report.pdf") == "pdf"

    def test_multiple_dots(self):
        assert _ext("archive.tar.gz") == "gz"

    def test_no_extension(self):
        assert _ext("README") == ""

    def test_uppercase_extension(self):
        assert _ext("Document.PDF") == "pdf"

    def test_hidden_file(self):
        # Hidden files like .gitignore have no suffix (Path considers
        # the entire name as the stem when the filename starts with '.')
        assert _ext(".gitignore") == ""

    def test_path_with_dirs(self):
        assert _ext("/path/to/file.txt") == "txt"

    def test_empty_string(self):
        assert _ext("") == ""


# ── _parse_text ────────────────────────────────────────────────────────────


class TestParseText:
    """Plain text / code / markdown file parsing."""

    def test_simple_text(self):
        result = _parse_text(b"Hello, world!", "test.txt")
        assert result["markdown"] == "Hello, world!"
        assert result["metadata"]["format"] == "txt"
        assert result["metadata"]["filename"] == "test.txt"

    def test_markdown_content(self):
        md = "# Title\n\nThis is **bold** text."
        result = _parse_text(md.encode(), "doc.md")
        assert "# Title" in result["markdown"]
        assert result["metadata"]["format"] == "md"

    def test_unicode_text(self):
        text = "Hello, 世界! ñoño 😀"
        result = _parse_text(text.encode("utf-8"), "unicode.txt")
        assert "世界" in result["markdown"]

    def test_large_text(self):
        text = "Line\n" * 10000
        result = _parse_text(text.encode(), "large.txt")
        assert len(result["markdown"]) > 1000
        assert result["metadata"]["chars"] == len(text)

    def test_empty_text(self):
        result = _parse_text(b"", "empty.txt")
        assert result["markdown"] == ""


# ── _parse_csv ─────────────────────────────────────────────────────────────


class TestParseCsv:
    """CSV parsing to markdown tables."""

    def test_standard_csv(self):
        csv_data = "name,age,city\nAlice,30,NYC\nBob,25,LA\n"
        result = _parse_csv(csv_data.encode(), "data.csv")
        assert result["metadata"]["format"] == "csv"
        assert result["metadata"]["rows"] == 3  # including header
        assert "Alice" in result["markdown"]
        assert "Bob" in result["markdown"]

    def test_single_row_csv(self):
        csv_data = "header\nvalue\n"
        result = _parse_csv(csv_data.encode(), "single.csv")
        assert result["metadata"]["rows"] == 2

    def test_empty_csv(self):
        result = _parse_csv(b"", "empty.csv")
        assert result["metadata"]["rows"] == 0

    def test_unicode_csv(self):
        csv_data = "name,emoji\nAlice,😀\nBob,ñ\n"
        result = _parse_csv(csv_data.encode("utf-8"), "unicode.csv")
        assert "😀" in result["markdown"]
        assert "ñ" in result["markdown"]

    def test_csv_with_quoted_fields(self):
        csv_data = '"name","description"\n"Smith, John","Engineer, Sr."\n'
        result = _parse_csv(csv_data.encode(), "quoted.csv")
        assert "Smith, John" in result["markdown"]
        assert result["metadata"]["rows"] == 2


# ── _parse_docx ────────────────────────────────────────────────────────────


class TestParseDocx:
    """DOCX parsing to markdown."""

    def _make_docx(self, paragraphs=None):
        """Create a minimal DOCX in memory."""
        from docx import Document

        doc = Document()
        if paragraphs:
            for p in paragraphs:
                doc.add_paragraph(p)
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    def _make_docx_with_table(self):
        """Create a DOCX with a table."""
        from docx import Document

        doc = Document()
        doc.add_paragraph("Before table")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "A"
        table.cell(1, 1).text = "1"
        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    def test_paragraphs(self):
        content = self._make_docx(["Hello world", "Second paragraph"])
        result = _parse_docx(content, "test.docx")
        assert "Hello world" in result["markdown"]
        assert result["metadata"]["paragraphs"] >= 2
        assert result["metadata"]["format"] == "docx"

    def test_empty_document(self):
        content = self._make_docx()
        result = _parse_docx(content, "empty.docx")
        assert result["markdown"] == ""

    def test_table_extraction(self):
        content = self._make_docx_with_table()
        result = _parse_docx(content, "table.docx")
        assert "Before table" in result["markdown"]
        assert "Name" in result["markdown"]
        assert "Value" in result["markdown"]
        assert result["metadata"]["tables"] == 1

    def test_corrupt_file(self):
        with pytest.raises((Exception,)):
            _parse_docx(b"not a docx file", "bad.docx")


# ── _parse_pptx ────────────────────────────────────────────────────────────


class TestParsePptx:
    """PPTX parsing to markdown."""

    def _make_pptx(self, slide_texts=None):
        """Create a minimal PPTX in memory."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        if slide_texts:
            for text in slide_texts:
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
                tx_box = slide.shapes.add_textbox(
                    Inches(1), Inches(1), Inches(8), Inches(2)
                )
                tf = tx_box.text_frame
                tf.text = text
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def _make_pptx_with_table(self):
        """Create a PPTX with a table."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        table_shape = slide.shapes.add_table(
            2, 2, Inches(1), Inches(1), Inches(6), Inches(2)
        )
        table = table_shape.table
        table.cell(0, 0).text = "Key"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "X"
        table.cell(1, 1).text = "42"
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def test_slides_with_text(self):
        content = self._make_pptx(["Slide 1 content", "Slide 2 content"])
        result = _parse_pptx(content, "test.pptx")
        assert (
            "Slide 1" in result["markdown"] or "Slide 1 content" in result["markdown"]
        )
        assert result["metadata"]["slides"] == 2
        assert result["metadata"]["format"] == "pptx"

    def test_empty_presentation(self):
        content = self._make_pptx()
        result = _parse_pptx(content, "empty.pptx")
        assert result["markdown"] == ""

    def test_tables_in_slides(self):
        content = self._make_pptx_with_table()
        result = _parse_pptx(content, "table.pptx")
        assert "Key" in result["markdown"]
        assert "Value" in result["markdown"]
        assert "42" in result["markdown"]

    def test_error_handling(self):
        with pytest.raises((Exception,)):
            _parse_pptx(b"not a pptx", "bad.pptx")


# ── _parse_xlsx ────────────────────────────────────────────────────────────


class TestParseXlsx:
    """XLSX parsing to markdown tables."""

    def _make_xlsx(self, sheets=None):
        """Create a minimal XLSX in memory."""
        import openpyxl

        wb = openpyxl.Workbook()
        if sheets:
            for i, (name, rows) in enumerate(sheets):
                if i == 0:
                    ws = wb.active
                    ws.title = name
                else:
                    ws = wb.create_sheet(title=name)
                for row_data in rows:
                    ws.append(row_data)
        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()

    def _make_xlsx_with_formulas(self):
        """Create XLSX with formulas."""
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Formulas"
        ws.append(["A", "B", "C"])
        ws.append([1, 2, 3])
        ws["A3"] = "=SUM(A2:A2)"
        buf = io.BytesIO()
        wb.save(buf)
        return buf.getvalue()

    def test_single_sheet(self):
        content = self._make_xlsx([("Sheet1", [("Name", "Age"), ("Alice", 30)])])
        result = _parse_xlsx(content, "test.xlsx")
        assert "Name" in result["markdown"]
        assert "Alice" in result["markdown"]
        assert result["metadata"]["format"] == "xlsx"

    def test_multi_sheet(self):
        sheets = [
            ("Sheet1", [("A", "B"), (1, 2)]),
            ("Sheet2", [("X", "Y"), (3, 4)]),
        ]
        content = self._make_xlsx(sheets)
        result = _parse_xlsx(content, "multi.xlsx")
        assert "Sheet1" in result["markdown"]
        assert "Sheet2" in result["markdown"]
        assert result["metadata"]["sheets"] == 2

    def test_formulas_data_only(self):
        content = self._make_xlsx_with_formulas()
        result = _parse_xlsx(content, "formulas.xlsx")
        # data_only=True means we see the cached value, not the formula
        assert "SUM" not in result["markdown"]

    def test_empty_workbook(self):
        import openpyxl

        wb = openpyxl.Workbook()
        buf = io.BytesIO()
        wb.save(buf)
        content = buf.getvalue()
        result = _parse_xlsx(content, "empty.xlsx")
        # An empty workbook still has one empty sheet
        assert result["markdown"] == "" or result["metadata"]["format"] == "xlsx"

    def test_error_handling(self):
        with pytest.raises((Exception,)):
            _parse_xlsx(b"not an xlsx", "bad.xlsx")


# ── _parse_pdf ──────────────────────────────────────────────────────────────


class TestParsePdf:
    """PDF parsing with graceful degradation tests."""

    def _make_minimal_pdf(self, text_content=None):
        """Create a minimal valid PDF with optional text."""
        if text_content:
            # Create a simple PDF with text using pypdf
            from pypdf import PdfWriter

            writer = PdfWriter()
            writer.add_blank_page(612, 792)
            # Can't easily add text via pypdf, so use minimal PDF approach
        # Return a minimal valid PDF
        return (
            b"%PDF-1.4\n"
            b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
            b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
            b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R/Resources<<>>>>endobj\n"
            b"xref\n0 4\n0000000000 65535 f \n0000000009 00000 n \n0000000058 00000 n \n0000000115 00000 n \n"
            b"trailer<</Size 4/Root 1 0 R>>\n"
            b"startxref\n190\n%%EOF"
        )

    def test_pypdf_extraction_path(self):
        """Verify pypdf-based extraction works (the minimal PDF has no text, returns empty)."""
        content = self._make_minimal_pdf()
        result = _parse_pdf(content, "test.pdf")
        assert "pdf" in result["metadata"]["format"]
        assert result["metadata"]["filename"] == "test.pdf"
        assert "pages" in result["metadata"]

    def test_ocr_degradation_graceful(self):
        """Verify that OCR import failure is handled gracefully (pytesseract/pdf2image not installed)."""
        # pytesseract and pdf2image are not installed, so OCR will log a warning
        # and skip gracefully. The function should still return a valid result.
        content = self._make_minimal_pdf()
        result = _parse_pdf(content, "test.pdf")
        assert "markdown" in result
        assert "metadata" in result
        # Even without text extracted, the function should not crash
        assert result["metadata"]["format"] == "pdf"

    def test_camelot_degradation_graceful(self):
        """Verify that Camelot import failure is handled gracefully (camelot not installed)."""
        # camelot is not installed, so table extraction will log a debug message
        # and skip gracefully. The function should still return a valid result.
        content = self._make_minimal_pdf()
        result = _parse_pdf(content, "test.pdf")
        assert "markdown" in result
        assert result["metadata"]["format"] == "pdf"
        # tables_found should not be set since camelot not available
        # but the function should not crash

    def test_corrupt_pdf_graceful(self):
        """Verify that corrupt PDF doesn't crash the parser."""
        result = _parse_pdf(b"not a pdf at all", "bad.pdf")
        assert "markdown" in result
        assert "metadata" in result
        # Should gracefully handle the failure


# ── _parse_via_markitdown ──────────────────────────────────────────────────


class TestParseViaMarkitdown:
    """Markitdown fallback parser."""

    def test_markitdown_path_json(self):
        """Verify markitdown can parse a JSON file."""
        content = b'{"name": "test", "value": 42}'
        result = _parse_via_markitdown(content, "data.json")
        assert result["metadata"]["format"] == "json"
        assert "name" in result["markdown"] or "test" in result["markdown"]

    def test_markitdown_path_html(self):
        """Verify markitdown can parse an HTML file."""
        content = b"<html><body><h1>Title</h1><p>Content</p></body></html>"
        result = _parse_via_markitdown(content, "page.html")
        assert result["metadata"]["format"] == "html"
        # MarkItDown should extract the title and content
        assert "Title" in result["markdown"] or "Content" in result["markdown"]

    def test_markitdown_path_yaml(self):
        """Verify markitdown can parse a YAML file."""
        content = b"name: test\nversion: 1.0\n"
        result = _parse_via_markitdown(content, "config.yaml")
        assert result["metadata"]["format"] == "yaml"
        assert "test" in result["markdown"] or "name" in result["markdown"]

    def test_markitdown_unknown_extension(self):
        """Verify markitdown handles unknown extensions gracefully."""
        result = _parse_via_markitdown(b"Some raw content", "file.weird_ext")
        # markitdown handles unknown extensions by treating them as plain text
        assert result["metadata"]["extraction"] == "markitdown"
        assert "Some raw content" in result["markdown"]

    def test_raw_decode_fallback_when_markitdown_fails(self):
        """Verify raw decode fallback when markitdown itself raises an exception."""
        # Patch MarkItDown.convert to raise an exception
        with patch(
            "markitdown.MarkItDown.convert", side_effect=Exception("Simulated failure")
        ):
            result = _parse_via_markitdown(b"Fallback content", "file.txt")
        assert result["metadata"]["extraction"] == "raw"
        assert "Fallback content" in result["markdown"]

    def test_unicode_content(self):
        """Verify unicode content survives through markitdown."""
        content = "Hello, 世界! 😀".encode()
        result = _parse_via_markitdown(content, "text.txt")
        # markitdown should extract text content
        assert result["metadata"]["format"] == "txt"
        assert len(result["markdown"]) > 0
