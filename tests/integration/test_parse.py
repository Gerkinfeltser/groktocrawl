"""Tests for parse-svc — document parsing endpoints.

Covers all supported formats, error cases, health, and metrics.
"""

import io

from fastapi.testclient import TestClient
from parse_svc.app import MAX_SIZE_MB, app

client = TestClient(app)


# ── Helper factories ────────────────────────────────────────────────────────


def _make_docx() -> bytes:
    """Create a minimal real DOCX in memory with paragraphs."""
    from docx import Document

    doc = Document()
    doc.add_paragraph("Hello from DOCX")
    doc.add_paragraph("Second paragraph")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx() -> bytes:
    """Create a minimal real PPTX in memory with one slide."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    tf = tx_box.text_frame
    tf.text = "Hello from PPTX"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_xlsx() -> bytes:
    """Create a minimal real XLSX in memory with one sheet."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Value"])
    ws.append(["Alice", 30])
    ws.append(["Bob", 25])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ── Health ──────────────────────────────────────────────────────────────────


def test_health():
    resp = client.get("/health")
    assert resp.status_code == 200
    data = resp.json()
    assert data["status"] == "ok"


# ── Metrics ─────────────────────────────────────────────────────────────────


def test_metrics():
    resp = client.get("/metrics")
    assert resp.status_code == 200
    assert "openmetrics-text" in resp.headers["content-type"]
    body = resp.text
    assert "# HELP" in body
    assert "# TYPE" in body
    assert "# EOF" in body.strip()


# ── Error cases ─────────────────────────────────────────────────────────────


def test_unsupported_format():
    resp = client.post(
        "/parse",
        files={"file": ("test.xyz", io.BytesIO(b"hello"), "application/octet-stream")},
    )
    assert resp.status_code == 400
    detail = resp.json()["detail"]
    assert "Unsupported format" in detail


def test_file_too_large():
    large_content = b"x" * (MAX_SIZE_MB * 1024 * 1024 + 1)
    resp = client.post(
        "/parse",
        files={"file": ("large.pdf", io.BytesIO(large_content), "application/pdf")},
    )
    assert resp.status_code == 413
    assert "too large" in resp.json()["detail"].lower()


def test_no_filename():
    """Upload without a filename → error.

    The code path returns 400, but the TestClient's multipart handling
    may return 422 when filename is empty.
    """
    resp = client.post(
        "/parse",
        files={"file": ("", io.BytesIO(b"content"), "text/plain")},
    )
    assert resp.status_code in (400, 422)


def test_no_extension():
    """Upload a file with no extension → 400."""
    resp = client.post(
        "/parse",
        files={"file": ("README", io.BytesIO(b"content"), "text/plain")},
    )
    assert resp.status_code == 400
    detail = resp.json()["detail"]
    assert "extension" in detail.lower()


# ── Supported formats ───────────────────────────────────────────────────────


def test_pdf_parsing():
    # Minimal valid PDF
    minimal_pdf = (
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
        b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R/Resources<<>>>>endobj\n"
        b"xref\n0 4\n0000000000 65535 f \n0000000009 00000 n \n0000000058 00000 n \n0000000115 00000 n \n"
        b"trailer<</Size 4/Root 1 0 R>>\n"
        b"startxref\n190\n%%EOF"
    )
    resp = client.post(
        "/parse",
        files={"file": ("test.pdf", io.BytesIO(minimal_pdf), "application/pdf")},
    )
    assert resp.status_code == 200
    data = resp.json()
    assert data["success"] is True
    assert data["data"]["metadata"]["format"] == "pdf"


def test_txt_parsing():
    resp = client.post(
        "/parse",
        files={"file": ("test.txt", io.BytesIO(b"hello world"), "text/plain")},
    )
    assert resp.status_code == 200
    data = resp.json()
    assert data["success"] is True
    assert "hello world" in data["data"]["markdown"]


def test_csv_parsing():
    csv_content = b"name,age\nAlice,30\nBob,25\n"
    resp = client.post(
        "/parse",
        files={"file": ("data.csv", io.BytesIO(csv_content), "text/csv")},
    )
    assert resp.status_code == 200
    data = resp.json()
    assert data["success"] is True
    assert data["data"]["metadata"]["format"] == "csv"
    assert "Alice" in data["data"]["markdown"]
    assert "Bob" in data["data"]["markdown"]


def test_docx_parsing():
    content = _make_docx()
    resp = client.post(
        "/parse",
        files={
            "file": (
                "report.docx",
                io.BytesIO(content),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        },
    )
    assert resp.status_code == 200
    data = resp.json()
    assert data["success"] is True
    assert data["data"]["metadata"]["format"] == "docx"
    assert "Hello from DOCX" in data["data"]["markdown"]


def test_pptx_parsing():
    content = _make_pptx()
    resp = client.post(
        "/parse",
        files={
            "file": (
                "slides.pptx",
                io.BytesIO(content),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert resp.status_code == 200
    data = resp.json()
    assert data["success"] is True
    assert data["data"]["metadata"]["format"] == "pptx"
    assert "Hello from PPTX" in data["data"]["markdown"]


def test_xlsx_parsing():
    content = _make_xlsx()
    resp = client.post(
        "/parse",
        files={
            "file": (
                "data.xlsx",
                io.BytesIO(content),
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            )
        },
    )
    assert resp.status_code == 200
    data = resp.json()
    assert data["success"] is True
    assert data["data"]["metadata"]["format"] == "xlsx"
    assert "Alice" in data["data"]["markdown"]
    assert "Bob" in data["data"]["markdown"]


def test_markdown_parsing():
    md_content = b"# Title\n\nThis is **bold** and *italic*."
    resp = client.post(
        "/parse",
        files={"file": ("doc.md", io.BytesIO(md_content), "text/markdown")},
    )
    assert resp.status_code == 200
    data = resp.json()
    assert data["success"] is True
    assert data["data"]["metadata"]["format"] == "md"
    assert "# Title" in data["data"]["markdown"]
    assert "bold" in data["data"]["markdown"]


def test_extension_detection():
    """DOCX with PK header — the parser detects extension first, then tries to parse."""
    resp = client.post(
        "/parse",
        files={
            "file": (
                "report.docx",
                io.BytesIO(b"PK\x03\x04fake"),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        },
    )
    # The extension is detected as docx; parsing may fail with 422
    assert resp.status_code in (200, 422)
